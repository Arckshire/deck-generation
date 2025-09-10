import io
import json
from datetime import datetime
from typing import List, Dict, Any, Optional

import pandas as pd
import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ---------- UI CONFIG ----------
st.set_page_config(
    page_title="EBR : KPI Analysis Deck",
    page_icon="📊",
    layout="wide",
)

PRIMARY = (23, 104, 172)  # default accent if template not providing a theme
LIGHT_BG = (245, 247, 250)

# ---------- HELPERS ----------
def load_table(file) -> Optional[pd.DataFrame]:
    name = file.name.lower()
    try:
        if name.endswith(".csv"):
            return pd.read_csv(file)
        if name.endswith(".xlsx") or name.endswith(".xls"):
            return pd.read_excel(file)
        if name.endswith(".json"):
            data = json.load(file)
            return pd.json_normalize(data)
    except Exception as e:
        st.warning(f"Could not read {file.name}: {e}")
        return None
    st.warning(f"Unsupported table file: {file.name}")
    return None

def safe_color(rgb_tuple):
    r, g, b = rgb_tuple
    return RGBColor(r, g, b)

def add_textbox(shape, title, subtitle=None, center=False, font_size=18, bold=False):
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = bold
    run.font.size = Pt(font_size)
    if center:
        p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = text_frame.add_paragraph()
        p2.text = subtitle
        p2.level = 1
        p2.font.size = Pt(max(12, font_size - 4))
        if center:
            p2.alignment = PP_ALIGN.CENTER

def add_kpi_card(slide, left_in, top_in, width_in, height_in, headline:str, value:str, note:str="", accent=PRIMARY):
    shape = slide.shapes.add_shape(
        autoshape_type_id=1,  # Rectangle
        left=Inches(left_in), top=Inches(top_in),
        width=Inches(width_in), height=Inches(height_in)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = safe_color((255, 255, 255))
    line = shape.line
    line.color.rgb = safe_color(accent)
    line.width = Pt(2)

    # Title
    tx_box = slide.shapes.add_textbox(
        Inches(left_in + 0.25), Inches(top_in + 0.2),
        Inches(width_in - 0.5), Inches(0.5)
    )
    add_textbox(tx_box, headline, center=False, font_size=16, bold=True)

    # Value
    val_box = slide.shapes.add_textbox(
        Inches(left_in + 0.25), Inches(top_in + 0.7),
        Inches(width_in - 0.5), Inches(0.8)
    )
    add_textbox(val_box, value, center=False, font_size=24, bold=True)

    # Note
    if note:
        note_box = slide.shapes.add_textbox(
            Inches(left_in + 0.25), Inches(top_in + 1.4),
            Inches(width_in - 0.5), Inches(height_in - 1.6)
        )
        add_textbox(note_box, note, center=False, font_size=12, bold=False)

def add_chart(slide, left, top, width, height, df:pd.DataFrame, x_col:str, y_col:str, chart_type:str="bar", series_name:str="Value"):
    chart_data = CategoryChartData()
    chart_data.categories = list(df[x_col].astype(str).values)
    chart_data.add_series(series_name, list(df[y_col].astype(float).values))

    if chart_type == "line":
        ctype = XL_CHART_TYPE.LINE_MARKERS
    else:
        ctype = XL_CHART_TYPE.COLUMN_CLUSTERED

    x, y, w, h = Inches(left), Inches(top), Inches(width), Inches(height)
    chart = slide.shapes.add_chart(ctype, x, y, w, h, chart_data).chart
    chart.has_legend = False
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.chart_area.format.fill.solid()
    chart.chart_area.format.fill.fore_color.rgb = safe_color((255,255,255))
    return chart

def add_screenshot(slide, img:Image.Image, left_in, top_in, width_in=None, height_in=None):
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(left_in), Inches(top_in),
                              width=Inches(width_in) if width_in else None,
                              height=Inches(height_in) if height_in else None)

def crop_image(img:Image.Image, bbox):
    return img.crop(bbox)

def add_speaker_notes(slide, notes:str):
    slide.notes_slide.notes_text_frame.text = notes

def create_presentation(base_template_bytes: Optional[bytes]) -> Presentation:
    if base_template_bytes:
        return Presentation(io.BytesIO(base_template_bytes))
    return Presentation()  # fallback to default theme

# ---------- SIDEBAR INPUTS ----------
st.sidebar.header("🔧 Inputs")
template_file = st.sidebar.file_uploader("Base PPTX template (optional)", type=["pptx"])
data_files = st.sidebar.file_uploader("Data files (CSV/XLSX/JSON)", type=["csv", "xlsx", "xls", "json"], accept_multiple_files=True)
screenshots = st.sidebar.file_uploader("Screenshots (PNG/JPG)", type=["png", "jpg", "jpeg"], accept_multiple_files=True)

st.sidebar.header("⚙️ Options")
title_text = st.sidebar.text_input("Slide Title", "Executive KPI Summary")
customer = st.sidebar.text_input("Customer", "")
date_range = st.sidebar.text_input("Date Range", "Last 3 months")
purpose = st.sidebar.selectbox("Purpose", ["C-level review", "CSM touchpoint", "Quarterly EBR", "Other"], index=0)
max_slides = st.sidebar.slider("Max slides", 1, 2, 2)
accent_color = st.sidebar.color_picker("Accent Color", "#1768AC")
include_speaker_notes = st.sidebar.checkbox("Include speaker notes", value=True)
enable_auto_crop = st.sidebar.checkbox("Auto-crop large screenshots into cards", value=True)

# ---------- MAIN ----------
st.title("📊 EBR : KPI Analysis Deck")
st.caption("Generate 1–2 adaptive slides with KPI cards, small charts, and screenshots.")

# Load tables
tables: List[pd.DataFrame] = []
if data_files:
    for f in data_files:
        df = load_table(f)
        if df is not None:
            tables.append(df)

# Data mapping (optional) to make a quick chart
st.subheader("Data Mapping (optional)")
st.write("Map a simple chart from one uploaded table.")
chart_df = None
x_col = y_col = None
chart_type = st.selectbox("Chart type", ["bar", "line"], index=0)
if tables:
    df_idx = st.selectbox("Select a table", list(range(len(tables))), format_func=lambda i: data_files[i].name)
    chart_df = tables[df_idx]
    x_col = st.selectbox("X axis", chart_df.columns.tolist())
    y_col = st.selectbox("Y axis (numeric)", chart_df.columns.tolist())
    # keep it simple: user must ensure numeric
else:
    st.info("Upload a CSV/XLSX/JSON to enable chart mapping (optional).")

# KPI Cards builder
st.subheader("Insight Cards")
st.write("Add KPI/insight cards that will be laid out across the slide(s).")
card_cols = st.columns([2, 2, 4])
with card_cols[0]:
    headline = st.text_input("Card headline (3–5 words)", "Tracking Rate")
with card_cols[1]:
    value = st.text_input("Card value", "90%")
with card_cols[2]:
    note = st.text_input("Short note (optional)", "Consistent improvement vs last month")

add_card = st.button("➕ Add card to queue")
if "cards" not in st.session_state:
    st.session_state.cards = []

if add_card:
    st.session_state.cards.append({"headline": headline, "value": value, "note": note})

if st.session_state.cards:
    st.write("Queued cards:")
    for i, c in enumerate(st.session_state.cards):
        st.write(f"- **{c['headline']}** — {c['value']} – {c['note']}")

# Screenshot cropping UI (basic)
st.subheader("Screenshots (optional)")
if screenshots:
    st.write("Uploaded screenshots:")
    for f in screenshots:
        st.image(f, caption=f.name, use_container_width=True)
else:
    st.caption("You can upload a single big screenshot; the app can auto-crop it into smaller cards (no OCR).")

# Generate button
if st.button("🚀 Generate PPTX"):
    # Prepare presentation
    accent_hex = accent_color.lstrip("#")
    accent_rgb = tuple(int(accent_hex[i:i+2], 16) for i in (0, 2, 4))
    prs = create_presentation(template_file.getvalue() if template_file else None)

    # Create 1 or 2 slides (blank layout index 6 is common, but templates vary)
    layout_idx = 6 if len(prs.slide_layouts) > 6 else 0

    # Decide how many slides
    total_cards = len(st.session_state.cards)
    total_imgs = len(screenshots) if screenshots else 0
    use_two_slides = (total_cards + total_imgs > 4) and (max_slides >= 2)

    def build_slide(title_str, cards_subset, imgs_subset, chart_cfg=None):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.6))
        add_textbox(title_box, title_str, font_size=28, bold=True)

        subtitle = []
        if customer: subtitle.append(f"Customer: {customer}")
        if date_range: subtitle.append(f"Range: {date_range}")
        if purpose: subtitle.append(f"Purpose: {purpose}")
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.4))
        add_textbox(sub_box, " • ".join(subtitle), font_size=12)

        # Card grid (up to 4 per slide)
        grid_positions = [
            (0.5, 1.5), (6.5, 1.5),
            (0.5, 4.0), (6.5, 4.0),
        ]
        card_w, card_h = 5.5, 2.2

        for idx, c in enumerate(cards_subset[:4]):
            left, top = grid_positions[idx]
            add_kpi_card(
                slide, left, top, card_w, card_h,
                headline=c["headline"], value=c["value"], note=c.get("note",""),
                accent=accent_rgb
            )

        # Chart (optional)
        if chart_cfg and chart_cfg.get("df") is not None:
            try:
                add_chart(
                    slide, left=0.5, top=6.3, width=11.5, height=2.3,
                    df=chart_cfg["df"], x_col=chart_cfg["x"], y_col=chart_cfg["y"],
                    chart_type=chart_cfg["type"], series_name=chart_cfg["series"]
                )
            except Exception as e:
                st.warning(f"Chart could not be added: {e}")

        # Screenshots (optional) — place below or next to cards
        img_left, img_top, img_w = 0.5, 6.3, 11.5
        placed = 0
        for img in imgs_subset[:1]:
            try:
                pil = Image.open(img).convert("RGB")
                add_screenshot(slide, pil, img_left, img_top, width_in=img_w)
                placed += 1
            except Exception as e:
                st.warning(f"Could not embed {getattr(img, 'name', 'image')}: {e}")

        # Speaker notes
        if include_speaker_notes:
            notes_lines = []
            for c in cards_subset[:4]:
                notes_lines.append(f"- {c['headline']}: {c['value']} — {c.get('note','')}")
            if chart_cfg and chart_cfg.get("df") is not None:
                notes_lines.append(f"- Chart: {chart_cfg['type']} of {chart_cfg['y']} by {chart_cfg['x']}")
            add_speaker_notes(slide, "\n".join(notes_lines))

    # Split assets across slides if needed
    cards1, cards2 = st.session_state.cards, []
    imgs1, imgs2 = screenshots or [], []
    if use_two_slides:
        half_cards = max(4, (total_cards + 1) // 2)
        cards1, cards2 = st.session_state.cards[:half_cards], st.session_state.cards[half_cards:]
        half_imgs = max(1, (total_imgs + 1) // 2)
        imgs1, imgs2 = imgs1[:half_imgs], imgs1[half_imgs:]

    # Chart config if mapped
    chart_cfg = None
    if chart_df is not None and x_col and y_col:
        chart_cfg = {"df": chart_df, "x": x_col, "y": y_col, "type": chart_type, "series": y_col}

    # Build slides
    build_slide(title_text, cards1, imgs1, chart_cfg)
    if use_two_slides:
        build_slide(f"{title_text} (cont.)", cards2, imgs2, chart_cfg=None)

    # Deliver file
    out = io.BytesIO()
    prs.save(out)
    st.success("PPTX generated successfully.")
    st.download_button(
        label="⬇️ Download EBR Deck (.pptx)",
        data=out.getvalue(),
        file_name=f"EBR_KPI_Deck_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

st.markdown("---")
st.caption("Tip: Add multiple cards for tracking %, milestone completeness, exceptions, etc. Map a quick chart from one table. Upload a base template to match your brand.")
